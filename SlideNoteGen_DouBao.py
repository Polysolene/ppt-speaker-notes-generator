"""
【全图片PPT 智能备注生成器】
前置：WPS/PowerPoint 另存为“图片格式PPT” (.pptx)
优势：保留原排版，无需本地Office截屏渲染。
"""

import openai
import base64
import io
import re
from os import getenv
from time import sleep
import logging
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ================= 豆包配置区 =================
DOUBAO_API_KEY = getenv("DOUBAO_API_KEY")
DOUBAO_BASE_URL = "https://ark.cn-beijing.volces.com/api/v3"
DOUBAO_MODEL = ["doubao-seed-2-0-pro-260215", "doubao-1-5-vision-pro-32k-250115"][1]

# 运行控制
RETRY_COUNT = 3
API_SLEEP_TIME = 1  # 豆包并发较高，1秒间隔足够（免费版建议保留1秒）
START_SLIDE = 1    # 默认从第1页开始
INITIAL_CONTEXT = "这是第一页，暂无上下文"  # 默认初始上下文
PPTX_PATH = "文档名.pptx"
PPTX_PATH_OUTPUT = PPTX_PATH
append_mode = True  # True 为追加模式，False 为覆盖模式

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)


def image_to_base64(img: Image.Image) -> str:
    buffer = io.BytesIO()
    img.save(buffer, format="PNG")
    return base64.b64encode(buffer.getvalue()).decode("utf-8")


def generate_speaker_note(client, img, memory_context):
    """适配豆包API的生成函数，内置重试逻辑"""
    img_base64 = image_to_base64(img)
    prompt = f"""
            学术辅助AI，请深度思考分析后提供与画面紧密结合的详略得当辅助解析

            <全局上下文>
            {memory_context}
            </全局上下文>
            
            动态字数控制：如果是目录、标题页、致谢页或直白的文字论述，请1句话带过，不废话。
            不重复字面内容/全局上下文中的内容，点出痛点、机制或隐藏逻辑。对复杂图表/结构仔细详尽分析；详解新出现的术语。
            若画面模糊难以识别，请保持诚实，不编造信息。
            严禁使用“这张幻灯片展示了”等凑字数的废话。开门见山。
            如果PPT内容非中文，翻译为中文。
            严禁出现任何LaTeX语法,可以使用PPT备注支持的上标/下标/特殊字符。
            如有多个分析点，请换行分段表达。避免一整段话论述。

            【输出格式】：
            严格按照以下 XML 标签输出，不要任何多余字符以及任何额外的 Markdown 代码块符号：

            <notes>
            （这里输出解析。）
            </notes>

            <context_update>
            （用精炼的语言，将本页的核心推进融入现有的主线逻辑中，供下一页参考。限 300 字以内。）
            </context_update>
            """

    for attempt in range(RETRY_COUNT):
        try:
            response = client.chat.completions.create(
                model=DOUBAO_MODEL,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{img_base64}"}}
                    ]
                }],
                temperature=0.1
            )
            return response.choices[0].message.content
        except Exception as e:
            logger.warning(f"豆包API第{attempt + 1}次重试，原因: {e}")
            sleep(5 * (attempt + 1))

    raise RuntimeError("豆包API熔断：重试次数耗尽")


def generate_contextual_notes(client, pptx_path, output_path, start_slide=START_SLIDE, initial_context=INITIAL_CONTEXT):
    """主循环：带进度保存与断点记忆"""
    prs = Presentation(pptx_path)
    memory_context = initial_context

    try:
        for i, slide in enumerate(prs.slides):
            slide_num = i + 1
            if slide_num < start_slide: continue

            logger.info(f"正在处理 [{slide_num}/{len(prs.slides)}]...")
            image_bytes = next((s.image.blob for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE), None)
            if not image_bytes:
                raise ValueError(f"第{slide_num}页未检测到图片，与预期不符！请检查PPT格式是否为图片版。")

            img = None
            try:
                img = Image.open(io.BytesIO(image_bytes))
                full_text = generate_speaker_note(client, img, memory_context)

                # 提取记忆
                match_ctx = re.search(r"<context_update>(.*?)</context_update>", full_text, re.DOTALL)
                memory_context = match_ctx.group(1).strip() if match_ctx else "逻辑延续中"

                # 提取备注
                match_notes = re.search(r"<notes>(.*?)</notes>", full_text, re.DOTALL)
                notes_text = match_notes.group(1).strip() if match_notes else full_text

                # 写入备注：支持追加或覆盖
                if append_mode and slide.has_notes_slide:
                    original_notes = slide.notes_slide.notes_text_frame.text
                    notes_text = f"{original_notes}\n\n---\n\n{notes_text}"  # 用分隔线区分

                # 写入
                slide.notes_slide.notes_text_frame.text = notes_text
                logger.info(f"-> 成功。记忆追踪（{len(memory_context)}字）：{memory_context.ljust(80)[:80]}...")

            except Exception as e:
                logger.error(f"🚨 中断于第 {slide_num} 页: {e}")
                print(f"\n续写建议上下文:\n{memory_context}\n")
                break
            finally:
                if img: img.close()

            sleep(API_SLEEP_TIME)
    finally:
        prs.save(output_path)
        logger.info(f"💾 进度已安全保存至: {output_path}")
        logger.info(f"\n【最终记忆】\n{memory_context}\n")


if __name__ == "__main__":
    doubao_client = openai.OpenAI(api_key=DOUBAO_API_KEY, base_url=DOUBAO_BASE_URL)
    generate_contextual_notes(doubao_client, PPTX_PATH, PPTX_PATH_OUTPUT)
